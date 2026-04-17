"""
Generate aicritic executive 1-pager slide (Current State / Future State / Benefits).
Run: python generate_slide.py
Output: aicritic_executive_slide.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Palette ────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
ACCENT     = RGBColor(0x00, 0xC2, 0xFF)   # electric blue
GREEN      = RGBColor(0x00, 0xE5, 0x96)   # mint green
AMBER      = RGBColor(0xFF, 0xA5, 0x00)   # amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xBE, 0xCC)
PANEL_BG   = RGBColor(0x16, 0x2A, 0x3E)   # slightly lighter navy for panels

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color, transparency=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=12, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=10.5, bullet="▸", color=WHITE, line_spacing=1.1):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def add_divider(slide, left, top, width, color=ACCENT, thickness=0.03):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(thickness)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def build_slide():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ── Full background ────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)

    # ── Header bar ────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 13.33, 1.05, PANEL_BG)
    add_divider(slide, 0, 1.05, 13.33, ACCENT, 0.04)

    add_text(slide, "aicritic", 0.35, 0.08, 3, 0.55,
             font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

    add_text(slide, "Multi-Model AI Code Review  ·  Executive Overview",
             3.2, 0.18, 7, 0.45, font_size=13, bold=False,
             color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    add_text(slide, "CONFIDENTIAL", 10.8, 0.22, 2.2, 0.4,
             font_size=8, bold=False, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

    # ── Section labels row ────────────────────────────────────────────────
    COL = [0.25, 4.56, 8.87]
    COL_W = 3.9

    labels = [
        ("CURRENT STATE", AMBER),
        ("FUTURE STATE",  GREEN),
        ("BENEFITS",      ACCENT),
    ]
    for (label, col_color), x in zip(labels, COL):
        add_rect(slide, x, 1.25, COL_W, 0.38, col_color)
        add_text(slide, label, x + 0.12, 1.27, COL_W - 0.2, 0.35,
                 font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.LEFT)

    # ── Panel backgrounds ─────────────────────────────────────────────────
    for x in COL:
        add_rect(slide, x, 1.63, COL_W, 5.55, PANEL_BG)

    # ── Current State content ─────────────────────────────────────────────
    add_text(slide, "Single-model reviews", COL[0] + 0.15, 1.72, 3.6, 0.35,
             font_size=11, bold=True, color=AMBER)
    add_bullet_box(slide, [
        "One LLM call — one perspective, one blind spot",
        "Developers manually review AI suggestions with no cross-check",
        "Security, coverage & migration reviews are manual or skipped",
        "No audit trail — no record of what was flagged or fixed",
        "Tool fragmentation: separate tools per concern (Snyk, SonarQube, manual PR review)",
        "Findings require engineer judgment to triage — no risk prioritisation",
    ], COL[0] + 0.15, 2.12, 3.65, 5.0, font_size=10, color=LIGHT_GREY)

    # ── Future State content ───────────────────────────────────────────────
    add_text(slide, "Three-model critic chain", COL[1] + 0.15, 1.72, 3.6, 0.35,
             font_size=11, bold=True, color=GREEN)
    add_bullet_box(slide, [
        "Claude Sonnet → primary analysis across 8 built-in review tools",
        "Gemini → independent cross-check, flags gaps & contradictions",
        "Claude Opus → final arbiter assigns risk levels (Low/Med/High/Critical)",
        "Optional Fixer stage applies approved fixes with full diff preview",
        "Works as a CLI today; available as @aicritic in VS Code Copilot Chat",
        "Role files control model, focus & strictness — no code changes needed",
    ], COL[1] + 0.15, 2.12, 3.65, 5.0, font_size=10, color=LIGHT_GREY)

    # ── Benefits content ───────────────────────────────────────────────────
    add_text(slide, "Higher-confidence, faster delivery", COL[2] + 0.15, 1.72, 3.6, 0.35,
             font_size=11, bold=True, color=ACCENT)
    add_bullet_box(slide, [
        "Catches issues a single model misses — adversarial cross-check by design",
        "Covers security, secrets, coverage, migrations, dependencies & more",
        "Runs on existing Copilot Enterprise licence — zero extra cost",
        "Automated fixes reduce engineer toil on mechanical remediation",
        "Consistent risk taxonomy across all teams and repositories",
        "Full markdown report + backup-before-write for audit & compliance",
    ], COL[2] + 0.15, 2.12, 3.65, 5.0, font_size=10, color=LIGHT_GREY)

    # ── Arrow connectors (text-based) ──────────────────────────────────────
    add_text(slide, "→", 4.2, 3.9, 0.5, 0.5,
             font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 8.5, 3.9, 0.5, 0.5,
             font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # ── Footer ─────────────────────────────────────────────────────────────
    add_divider(slide, 0, 7.18, 13.33, PANEL_BG, 0.04)
    add_rect(slide, 0, 7.22, 13.33, 0.28, PANEL_BG)

    footer_items = [
        ("Phase 1 — CLI  ✓ Done", 0.35, LIGHT_GREY),
        ("Phase 2 — Copilot Extension  ✓ Done", 3.8, LIGHT_GREY),
        ("Phase 3 — Internal Hosting  ◌ Planned", 8.0, AMBER),
    ]
    for text, x, color in footer_items:
        add_text(slide, text, x, 7.23, 3.8, 0.25,
                 font_size=8, color=color, align=PP_ALIGN.LEFT)

    out = "aicritic_executive_slide.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build_slide()
